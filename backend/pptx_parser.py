import io
import re
from lxml import etree
from pptx import Presentation


def _extract_all_text_from_xml(element) -> list[str]:
    """Recursively extract ALL text from any XML element, including SmartArt,
    grouped shapes, and nested structures that python-pptx might miss."""
    texts = []
    # Namespace map for PowerPoint XML
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    # Find ALL text runs in the XML tree (catches SmartArt, grouped shapes, etc.)
    for t_elem in element.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
        text = t_elem.text
        if text and text.strip():
            texts.append(text.strip())

    return texts


def _extract_from_shape_recursive(shape, depth=0) -> list[str]:
    """Recursively extract text from a shape and all its children."""
    texts = []

    # Extract from text frames (standard text boxes, titles, etc.)
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            # Collect all runs in the paragraph to preserve full sentences
            para_text = paragraph.text.strip()
            if para_text:
                texts.append(para_text)

    # Extract from tables
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                texts.append(" | ".join(row_texts))

    # Recurse into grouped shapes (any depth)
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        try:
            for child_shape in shape.shapes:
                texts.extend(_extract_from_shape_recursive(child_shape, depth + 1))
        except Exception:
            pass

    # Fallback: parse the raw XML for anything we might have missed
    # This catches SmartArt, diagrams, and other complex structures
    try:
        xml_element = shape._element
        xml_texts = _extract_all_text_from_xml(xml_element)
        # Only add texts not already captured via the text_frame
        existing = set(texts)
        for t in xml_texts:
            if t not in existing:
                texts.append(t)
                existing.add(t)
    except Exception:
        pass

    return texts


def extract_text_from_pptx(file_path_or_bytes) -> str:
    """Extract ALL text content from a PPTX file, including complex structures.

    Handles: text boxes, titles, tables, grouped shapes (nested),
    SmartArt, diagrams, and any other text-containing elements.

    Args:
        file_path_or_bytes: Either a file path (str) or raw bytes of the PPTX file.

    Returns:
        Concatenated text from all slides, with slide separators.
    """
    if isinstance(file_path_or_bytes, bytes):
        prs = Presentation(io.BytesIO(file_path_or_bytes))
    else:
        prs = Presentation(file_path_or_bytes)

    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        # Extract from all shapes recursively
        for shape in slide.shapes:
            shape_texts = _extract_from_shape_recursive(shape)
            slide_texts.extend(shape_texts)

        # Also extract from slide notes (often contains extra CV info)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    slide_texts.append(f"[Note] {text}")

        if slide_texts:
            all_text.append(f"\n--- Slide {slide_num} ---")
            all_text.extend(slide_texts)

    result = "\n".join(all_text)

    # Clean up excessive whitespace while preserving structure
    result = re.sub(r'\n{3,}', '\n\n', result)

    # Fix mojibake: if UTF-8 text was decoded as Windows-1252, re-encode properly
    # (PPTX from Windows often produces cp1252 mojibake like "SantÃ©" instead of "Santé")
    try:
        fixed = result.encode('cp1252').decode('utf-8')
        result = fixed
    except (UnicodeDecodeError, UnicodeEncodeError):
        pass  # Not mojibake, keep original

    return result
